#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CDC健康检测系统 - 生理特征计算PPT（2页）
介绍HR→Mi→Tcr递推算法
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色
C_DARK = RGBColor(15, 25, 50)
C_BLUE = RGBColor(0, 120, 255)
C_CYAN = RGBColor(0, 220, 255)
C_ORANGE = RGBColor(255, 140, 0)
C_RED = RGBColor(255, 60, 60)
C_PURPLE = RGBColor(150, 80, 255)
C_GREEN = RGBColor(0, 200, 120)
C_WHITE = RGBColor(255, 255, 255)
C_LIGHT = RGBColor(248, 250, 255)
C_GRAY = RGBColor(100, 110, 120)

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_card(slide, left, top, width, height, fill=None, line_color=None):
    if fill is None:
        fill = C_WHITE
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if line_color:
        card.line.color.rgb = line_color
    else:
        card.line.fill.background()
    return card

# ========== 第1页：计算原理 ==========
def slide1_principle():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_DARK)
    
    # 装饰圆
    for x, y, size in [(Inches(-2), Inches(-1), Inches(5)), (Inches(10), Inches(5), Inches(5))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BLUE
        c.line.fill.background()
    
    # 标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心算法：HR → Mi → Tcr 递推模型"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(10), Inches(0.4))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "CORE ALGORITHM: HEART RATE → METABOLIC RATE → CORE TEMPERATURE"
    p.font.size = Pt(12)
    p.font.color.rgb = C_CYAN
    
    # 流程图
    # HR → Mi → Tcr
    boxes = [
        ("❤️", "心率 HR", "实时采集\n蓝牙心率带", Inches(0.8), C_RED),
        ("→", "", "", Inches(3.8), C_ORANGE),
        ("🔥", "代谢率 Mi", "能量消耗\n劳动强度", Inches(4.8), C_ORANGE),
        ("→", "", "", Inches(7.8), C_ORANGE),
        ("🌡️", "核心体温 Tcr", "预测值\n热风险评估", Inches(8.8), C_BLUE),
    ]
    
    for icon, name, desc, x, color in boxes:
        if icon == "→":
            # 箭头
            arrow = s.shapes.add_textbox(x, Inches(2.5), Inches(1), Inches(1))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(48)
            p.font.color.rgb = color
        else:
            # 卡片
            card = add_card(s, x, Inches(2), Inches(3.5), Inches(2.5), RGBColor(20, 40, 80), color)
            
            # 图标
            icon_box = s.shapes.add_textbox(x, Inches(2.2), Inches(3.5), Inches(1.2))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(56)
            p.alignment = PP_ALIGN.CENTER
            
            # 名称
            name_box = s.shapes.add_textbox(x + 0.2, Inches(3.5), Inches(3.1), Inches(0.5))
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            desc_box = s.shapes.add_textbox(x + 0.2, Inches(4.1), Inches(3.1), Inches(0.8))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(180, 200, 220)
            p.alignment = PP_ALIGN.CENTER
    
    # 公式卡片
    formula_bg = add_card(s, Inches(0.5), Inches(5), Inches(12.3), Inches(2.2), RGBColor(10, 30, 60))
    
    # 公式标题
    f_title = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(5), Inches(0.5))
    tf = f_title.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 核心计算公式"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    
    # Mi公式
    mi_box = s.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.6))
    tf = mi_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Mi = 65 + (HR - HRrest) / (180 - 0.65×Age - HRrest) × [(41.7 - 0.22×Age) × W^(2/3) - 65]"
    p.font.size = Pt(14)
    p.font.color.rgb = C_CYAN
    
    # Mi说明
    mi_desc = s.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.4))
    tf = mi_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "劳动代谢率 = 65 + 根据心率与静息心率差值 × 人体表面积参数"
    p.font.size = Pt(11)
    p.font.color.rgb = C_GRAY
    
    # Tcr公式
    tcr_box = s.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.4))
    tf = tcr_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Tcr(t+1) = Tcr(t) + 0.0036 × (Mi - 55) × 0.0952"
    p.font.size = Pt(14)
    p.font.color.rgb = C_CYAN

# ========== 第2页：创新优势 ==========
def slide2_advantages():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_LIGHT)
    
    # 标题栏
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = C_DARK
    header.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "算法创新优势"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    page = s.shapes.add_textbox(Inches(12), Inches(0.25), Inches(1), Inches(0.5))
    tf = page.text_frame
    p = tf.paragraphs[0]
    p.text = "02"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 四大优势
    advantages = [
        ("⏰", "提前预警", "15-30分钟", "通过心率变化\n提前预测体温上升\n及早发现热射病风险", C_RED),
        ("📊", "多参数融合", "综合分析", "融合心率、年龄、体重\n环境温湿度\n服装热阻等参数", C_BLUE),
        ("🎯", "实时递推", "秒级更新", "基于前一时刻Tcr\n递推计算下一时刻\n持续跟踪体温变化", C_ORANGE),
        ("🔬", "科学依据", "国际标准", "基于PMV-PPD模型\n和生理学研究成果\n计算结果可靠", C_PURPLE),
    ]
    
    positions = [(0.5, 1.2), (3.5, 1.2), (6.5, 1.2), (9.5, 1.2)]
    
    for i, (icon, title_text, tag, desc, color) in enumerate(advantages):
        x, y = positions[i]
        
        # 卡片
        card = add_card(s, Inches(x), Inches(y), Inches(2.9), Inches(3.5), C_WHITE, color)
        
        # 顶部色条
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.9), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # 图标
        icon_box = s.shapes.add_textbox(Inches(x), Inches(y + 0.3), Inches(2.9), Inches(1))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.4), Inches(2.5), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER
        
        # 标签
        tag_bg = add_card(s, Inches(x + 0.5), Inches(y + 2), Inches(1.9), Inches(0.4), RGBColor(245, 245, 245))
        tag_box = s.shapes.add_textbox(Inches(x + 0.5), Inches(y + 2.08), Inches(1.9), Inches(0.35))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        d_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 2.5), Inches(2.5), Inches(1))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = C_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # 底部对比
    compare_bg = add_card(s, Inches(0.5), Inches(5), Inches(12.3), Inches(2.2), C_WHITE)
    
    # 对比标题
    c_title = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(4), Inches(0.5))
    tf = c_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🔄 传统 vs 我们的方案"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    
    # 传统方式
    trad_box = s.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(5.5), Inches(1.2))
    tf = trad_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "❌ 传统方式"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_RED
    p = tf.add_paragraph()
    p.text = "• 仅测量皮肤表面温度\n• 依赖人工定期检测\n• 发现时往往已发病"
    p.font.size = Pt(12)
    p.font.color.rgb = C_GRAY
    
    # 我们的方案
    our_box = s.shapes.add_textbox(Inches(6.8), Inches(5.8), Inches(5.8), Inches(1.2))
    tf = our_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✅ 我们的方案"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_GREEN
    p = tf.add_paragraph()
    p.text = "• 通过心率预测核心体温\n• 实时自动监测\n• 提前15-30分钟预警"
    p.font.size = Pt(12)
    p.font.color.rgb = C_GRAY

# ========== 生成 ==========
print("🎨 正在生成生理特征计算PPT...")

slide1_principle()
slide2_advantages()

output_path = "/workspace/projects/CDC健康检测_生理算法PPT.pptx"
prs.save(output_path)
print(f"✅ PPT生成成功！")
print(f"📁 文件：{output_path}")
print(f"📊 共2页 - HR→Mi→Tcr递推算法介绍")
