#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CDC健康检测系统 - 产品介绍PPT（3页）
简洁专业的参赛作品展示
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色
C_DARK = RGBColor(15, 25, 50)
C_BLUE = RGBColor(0, 120, 255)
C_CYAN = RGBColor(0, 220, 255)
C_ORANGE = RGBColor(255, 140, 0)
C_RED = RGBColor(255, 60, 60)
C_GREEN = RGBColor(0, 200, 120)
C_PURPLE = RGBColor(150, 80, 255)
C_WHITE = RGBColor(255, 255, 255)
C_LIGHT = RGBColor(248, 250, 255)
C_GRAY = RGBColor(100, 110, 120)

IMG_DIR = "/workspace/projects/ppt_images"

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_card(slide, left, top, width, height, fill=None, line_color=None):
    if fill is None:
        fill = C_WHITE
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if line_color:
        card.line.color.rgb = line_color
    else:
        card.line.fill.background()
    return card

def add_image(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)

# ========== 第1页：封面 ==========
def slide1_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_DARK)
    
    # 大装饰圆
    for x, y, size in [(Inches(-2), Inches(-2), Inches(6)), (Inches(9), Inches(4), Inches(6))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BLUE
        c.line.fill.background()
    
    # 小装饰圆
    for x, y, size in [(Inches(11), Inches(-1), Inches(3)), (Inches(-1), Inches(5), Inches(2))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0, 80, 160)
        c.line.fill.background()
    
    # 比赛标签
    tag = add_card(s, Inches(3.5), Inches(1.2), Inches(6.5), Inches(0.7), C_ORANGE)
    t = s.shapes.add_textbox(Inches(3.5), Inches(1.35), Inches(6.5), Inches(0.5))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "🏆 计算机设计大赛 · Web应用与开发赛道"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 主标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.5), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "CDC健康检测系统"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub = s.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.5), Inches(0.6))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "基于现代Web技术的热舒适度实时监测平台"
    p.font.size = Pt(26)
    p.font.color.rgb = C_CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # 功能标签
    features = ["心率监测", "代谢计算", "体温预测", "CDC评估", "智能预警"]
    x = 1.5
    for f in features:
        tag = add_card(s, Inches(x), Inches(4.6), Inches(2), Inches(0.5), RGBColor(0, 100, 180))
        t = s.shapes.add_textbox(Inches(x), Inches(4.7), Inches(2), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = f
        p.font.size = Pt(14)
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        x += 2.3
    
    # 技术栈
    techs = ["Next.js 16", "React 19", "Tailwind CSS", "Supabase", "TypeScript"]
    x = 1.5
    for tech in techs:
        tag = add_card(s, Inches(x), Inches(5.3), Inches(2), Inches(0.5), RGBColor(30, 50, 80), C_BLUE)
        t = s.shapes.add_textbox(Inches(x), Inches(5.4), Inches(2), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(12)
        p.font.color.rgb = C_CYAN
        p.alignment = PP_ALIGN.CENTER
        x += 2.3
    
    # 底部条
    bottom = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7), prs.slide_width, Inches(0.5))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = C_ORANGE
    bottom.line.fill.background()
    
    # 底部文字
    bt = s.shapes.add_textbox(0, Inches(7.05), prs.slide_width, Inches(0.4))
    tf = bt.text_frame
    p = tf.paragraphs[0]
    p.text = "实时监测 · 智能预警 · 健康保障"
    p.font.size = Pt(16)
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# ========== 第2页：核心功能与创新 ==========
def slide2_features():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_LIGHT)
    
    # 标题栏
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = C_DARK
    header.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心功能与技术创新"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    # 左侧：核心功能
    func_title = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3), Inches(0.5))
    tf = func_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 核心功能"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    
    functions = [
        ("❤️", "心率监测", "蓝牙心率带实时采集", C_RED),
        ("🔥", "代谢计算", "HR→Mi→Tcr递推算法", C_ORANGE),
        ("🌡️", "体温预测", "提前预警热射病", C_BLUE),
        ("📊", "CDC评估", "PMV-PPD热舒适度", C_PURPLE),
        ("👔", "服装建议", "基于环境智能推荐", C_GREEN),
    ]
    
    y = 1.8
    for icon, title_text, desc, color in functions:
        card = add_card(s, Inches(0.5), Inches(y), Inches(5.5), Inches(0.9), C_WHITE, color)
        icon_box = s.shapes.add_textbox(Inches(0.7), Inches(y + 0.15), Inches(0.7), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(24)
        t_box = s.shapes.add_textbox(Inches(1.5), Inches(y + 0.15), Inches(2), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        d_box = s.shapes.add_textbox(Inches(1.5), Inches(y + 0.5), Inches(4.3), Inches(0.4))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = C_GRAY
        y += 1.0
    
    # 右侧：技术创新
    tech_title = s.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(4), Inches(0.5))
    tf = tech_title.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ 技术创新"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    
    innovations = [
        ("🚀", "HR→Mi→Tcr递推算法", "通过心率预测核心体温"),
        ("🔄", "多参数融合模型", "综合环境+生理多维数据"),
        ("☁️", "云端实时同步", "数据云端上传多端同步"),
        ("📱", "响应式设计", "PC/平板/手机自适应"),
    ]
    
    y = 1.8
    for icon, title_text, desc in innovations:
        card = add_card(s, Inches(6.5), Inches(y), Inches(6.3), Inches(1.15), C_WHITE, C_BLUE)
        icon_box = s.shapes.add_textbox(Inches(6.7), Inches(y + 0.25), Inches(0.7), Inches(0.7))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        t_box = s.shapes.add_textbox(Inches(7.5), Inches(y + 0.15), Inches(5), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        d_box = s.shapes.add_textbox(Inches(7.5), Inches(y + 0.55), Inches(5), Inches(0.5))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = C_GRAY
        y += 1.25

# ========== 第3页：技术架构与总结 ==========
def slide3_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_DARK)
    
    # 装饰圆
    for x, y, size in [(Inches(-1), Inches(-1), Inches(4)), (Inches(10), Inches(5), Inches(4))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BLUE
        c.line.fill.background()
    
    # 标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "技术架构与参赛信息"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    # 左侧：技术架构
    arch_title = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(3), Inches(0.5))
    tf = arch_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🛠️ 技术架构"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    
    layers = [
        ("🌐", "用户端", "Web浏览器 + 响应式设计"),
        ("⚙️", "前端", "Next.js 16 + React 19"),
        ("🔌", "API", "Next.js API Routes"),
        ("☁️", "数据库", "Supabase PostgreSQL"),
        ("📡", "设备", "Web Bluetooth"),
    ]
    
    y = 2.0
    for icon, name, desc in layers:
        card = add_card(s, Inches(0.5), Inches(y), Inches(5.5), Inches(0.8), RGBColor(20, 40, 80))
        icon_box = s.shapes.add_textbox(Inches(0.7), Inches(y + 0.15), Inches(0.6), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(22)
        n_box = s.shapes.add_textbox(Inches(1.4), Inches(y + 0.15), Inches(1.5), Inches(0.4))
        tf = n_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_CYAN
        d_box = s.shapes.add_textbox(Inches(1.4), Inches(y + 0.45), Inches(4.4), Inches(0.4))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(180, 200, 220)
        y += 0.88
    
    # 右侧：参赛信息
    comp_title = s.shapes.add_textbox(Inches(6.5), Inches(1.4), Inches(4), Inches(0.5))
    tf = comp_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🏆 参赛信息"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    
    comp_card = add_card(s, Inches(6.5), Inches(2.0), Inches(6.3), Inches(2.6), RGBColor(20, 40, 80), C_ORANGE)
    
    info = [
        ("📋", "赛事", "第十五届大学生计算机设计大赛"),
        ("🖥️", "赛道", "Web应用与开发"),
        ("💻", "技术", "Next.js + React + TypeScript"),
        ("☁️", "服务", "Supabase 云数据库"),
    ]
    
    y = 2.2
    for icon, label, content in info:
        l_box = s.shapes.add_textbox(Inches(6.8), Inches(y), Inches(1.2), Inches(0.4))
        tf = l_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = C_ORANGE
        c_box = s.shapes.add_textbox(Inches(8), Inches(y), Inches(4.5), Inches(0.4))
        tf = c_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon + " " + content
        p.font.size = Pt(13)
        p.font.color.rgb = C_WHITE
        y += 0.5
    
    # 应用场景
    scene_title = s.shapes.add_textbox(Inches(6.5), Inches(4.8), Inches(4), Inches(0.5))
    tf = scene_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🌍 应用场景"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    
    scenes = [
        ("🏗️", "高温作业"),
        ("🏃", "运动训练"),
        ("🏥", "医疗健康"),
    ]
    
    x = 6.5
    for icon, name in scenes:
        card = add_card(s, Inches(x), Inches(5.5), Inches(2), Inches(0.9), RGBColor(20, 40, 80))
        icon_box = s.shapes.add_textbox(Inches(x), Inches(5.6), Inches(2), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
        n_box = s.shapes.add_textbox(Inches(x), Inches(6.1), Inches(2), Inches(0.4))
        tf = n_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        x += 2.2
    
    # 感谢
    thanks = s.shapes.add_textbox(0, Inches(6.8), prs.slide_width, Inches(0.6))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听！欢迎交流！"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# ========== 生成 ==========
print("🎨 正在生成产品介绍PPT（3页）...")

slide1_cover()
slide2_features()
slide3_summary()

output_path = "/workspace/projects/CDC健康检测_产品介绍PPT.pptx"
prs.save(output_path)
print(f"✅ PPT生成成功！")
print(f"📁 文件：{output_path}")
print(f"📊 共3页 - 封面 + 功能创新 + 技术总结")
