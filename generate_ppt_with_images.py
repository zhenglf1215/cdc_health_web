#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CDC健康检测系统 - Web开发大赛PPT（多图版）
聚焦：技术架构 + 开发创新点 + 前端亮点
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 16:9比例
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案
C_DARK = RGBColor(15, 25, 50)
C_BLUE = RGBColor(0, 120, 255)
C_CYAN = RGBColor(0, 220, 255)
C_ORANGE = RGBColor(255, 140, 0)
C_PURPLE = RGBColor(150, 80, 255)
C_GREEN = RGBColor(0, 200, 120)
C_WHITE = RGBColor(255, 255, 255)
C_LIGHT = RGBColor(248, 250, 255)
C_GRAY = RGBColor(100, 110, 120)

IMG_DIR = "/workspace/projects/ppt_images"

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_card(slide, left, top, width, height, fill=None, line_color=None):
    if fill is None:
        fill = C_WHITE
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if line_color:
        card.line.color.rgb = line_color
    else:
        card.line.fill.background()
    return card

def add_image(slide, path, left, top, width, height):
    """添加图片"""
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)

# ========== 第1页：封面 ==========
def slide1_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景图
    bg_path = os.path.join(IMG_DIR, "cover.jpg")
    if os.path.exists(bg_path):
        add_bg(s, C_DARK)
        # 添加半透明覆盖层
        overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(15, 25, 50)
        overlay.line.fill.background()
        # 图片放底部
        add_image(s, bg_path, Inches(6), Inches(3), Inches(7), Inches(5))
    else:
        add_bg(s, C_DARK)
        # 装饰圆
        for x, y, size in [(Inches(-1), Inches(-2), Inches(5)), 
                           (Inches(9), Inches(4), Inches(6))]:
            c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
            c.fill.solid()
            c.fill.fore_color.rgb = C_BLUE
            c.line.fill.background()
    
    # 比赛标签
    tag = add_card(s, Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.7), C_ORANGE)
    t = s.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(4.5), Inches(0.5))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "🏆 计算机设计大赛 Web赛道"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 主标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(8), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "CDC健康检测系统"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    # 副标题
    sub = s.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(8), Inches(0.6))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "现代Web前端技术开发实践"
    p.font.size = Pt(24)
    p.font.color.rgb = C_CYAN
    
    # 技术栈标签
    techs = ["Next.js 16", "React 19", "Tailwind CSS", "Supabase", "TypeScript"]
    x = 0.5
    for tech in techs:
        tag = add_card(s, Inches(x), Inches(4.5), Inches(2.2), Inches(0.5), RGBColor(0, 100, 180))
        t = s.shapes.add_textbox(Inches(x), Inches(4.6), Inches(2.2), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(13)
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        x += 2.5
    
    # 底部
    bottom = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7), prs.slide_width, Inches(0.5))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = C_ORANGE
    bottom.line.fill.background()

# ========== 第2页：技术架构 ==========
def slide2_architecture():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_LIGHT)
    
    # 架构图背景
    arch_path = os.path.join(IMG_DIR, "architecture.jpg")
    if os.path.exists(arch_path):
        add_image(s, arch_path, Inches(6.5), Inches(1), Inches(6.5), Inches(5.5))
    
    # 标题栏
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = C_DARK
    header.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "技术架构"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    page = s.shapes.add_textbox(Inches(12), Inches(0.25), Inches(1), Inches(0.5))
    tf = page.text_frame
    p = tf.paragraphs[0]
    p.text = "01"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 左侧架构说明
    layers = [
        ("🌐", "用户端", "Web浏览器 / 响应式设计", C_BLUE),
        ("⚙️", "前端框架", "Next.js 16 + React 19", C_PURPLE),
        ("🔌", "API层", "Next.js API Routes", C_CYAN),
        ("☁️", "云服务", "Supabase PostgreSQL", C_GREEN),
        ("📡", "设备层", "Web Bluetooth", C_ORANGE),
    ]
    
    y = 1.3
    for icon, name, desc, color in layers:
        # 图标
        icon_box = s.shapes.add_textbox(Inches(0.3), Inches(y), Inches(0.8), Inches(0.9))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        
        # 名称
        name_box = s.shapes.add_textbox(Inches(1.2), Inches(y + 0.1), Inches(2), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        
        # 描述
        desc_box = s.shapes.add_textbox(Inches(1.2), Inches(y + 0.5), Inches(4), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = C_GRAY
        
        # 连接线
        if y < 5.5:
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y + 0.95), Inches(0.05), Inches(0.2))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(200, 200, 200)
            line.line.fill.background()
        
        y += 1.15

# ========== 第3页：前端亮点 + 技术栈图 ==========
def slide3_highlights():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_DARK)
    
    # 技术栈图
    tech_path = os.path.join(IMG_DIR, "tech_stack.jpg")
    if os.path.exists(tech_path):
        add_image(s, tech_path, Inches(7.5), Inches(1.5), Inches(5.5), Inches(4))
    
    # 装饰
    for x, y, size in [(Inches(-1), Inches(-1), Inches(3))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0, 80, 160)
        c.line.fill.background()
    
    # 标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "前端开发亮点"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(1), Inches(6), Inches(0.4))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "FRONTEND DEVELOPMENT HIGHLIGHTS"
    p.font.size = Pt(12)
    p.font.color.rgb = C_CYAN
    
    page = s.shapes.add_textbox(Inches(12), Inches(0.35), Inches(1), Inches(0.5))
    tf = page.text_frame
    p = tf.paragraphs[0]
    p.text = "02"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 六大亮点 - 左侧
    highlights = [
        ("⚡", "Next.js 16", "App Router服务端渲染\nTypeScript类型安全\nSEO友好", C_BLUE),
        ("🎨", "Tailwind CSS", "原子化CSS\n响应式设计\n主题定制", C_PURPLE),
        ("🧩", "shadcn/ui", "Radix UI组件\n无障碍支持\n可定制主题", C_CYAN),
        ("📊", "数据可视化", "实时图表\n交互体验\n数据筛选", C_GREEN),
    ]
    
    positions = [(0.4, 1.5), (0.4, 3.5), (0.4, 5.5),
                (4.2, 1.5), (4.2, 3.5), (4.2, 5.5)]
    
    for i, (icon, title_text, desc, color) in enumerate(highlights):
        x, y = positions[i]
        
        # 卡片
        card = add_card(s, Inches(x), Inches(y), Inches(3.5), Inches(1.8), RGBColor(20, 40, 80), color)
        
        # 图标
        icon_box = s.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(0.8), Inches(0.8))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(32)
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(x + 0.9), Inches(y + 0.2), Inches(2.4), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
        # 描述
        d_box = s.shapes.add_textbox(Inches(x + 0.9), Inches(y + 0.7), Inches(2.4), Inches(1))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(180, 200, 220)

# ========== 第4页：功能展示 + 仪表盘图 ==========
def slide4_features():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_LIGHT)
    
    # 仪表盘图
    dash_path = os.path.join(IMG_DIR, "dashboard.jpg")
    if os.path.exists(dash_path):
        add_image(s, dash_path, Inches(6.8), Inches(1.2), Inches(6.2), Inches(4.5))
    
    # 标题栏
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = C_DARK
    header.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "功能展示"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    page = s.shapes.add_textbox(Inches(12), Inches(0.25), Inches(1), Inches(0.5))
    tf = page.text_frame
    p = tf.paragraphs[0]
    p.text = "03"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 左侧功能说明
    features = [
        ("👤", "用户管理", "登录注册 · 角色权限 · 头像系统", C_BLUE),
        ("📈", "数据看板", "实时心率 · 趋势图表 · 数据筛选", C_GREEN),
        ("🌡️", "体征监测", "多参数显示 · 实时更新 · 预警提醒", C_ORANGE),
        ("📱", "响应式设计", "PC/平板/手机 · 自适应布局", C_PURPLE),
        ("🔄", "状态管理", "React Context · 数据持久化", C_CYAN),
        ("🔒", "安全设计", "RLS权限 · 数据校验 · XSS防护", RGBColor(255, 80, 120)),
    ]
    
    y = 1.3
    for icon, title_text, desc, color in features:
        # 图标
        icon_box = s.shapes.add_textbox(Inches(0.3), Inches(y), Inches(0.8), Inches(0.9))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(24)
        
        # 色块
        color_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(y + 0.25), Inches(0.08), Inches(0.4))
        color_bar.fill.solid()
        color_bar.fill.fore_color.rgb = color
        color_bar.line.fill.background()
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(1.4), Inches(y + 0.15), Inches(5), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        
        # 描述
        d_box = s.shapes.add_textbox(Inches(1.4), Inches(y + 0.55), Inches(5), Inches(0.4))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = C_GRAY
        
        y += 1.0

# ========== 第5页：总结 ==========
def slide5_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_DARK)
    
    # 装饰
    for x, y, size in [(Inches(-1), Inches(-1), Inches(4)), (Inches(10), Inches(5), Inches(5))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BLUE
        c.line.fill.background()
    
    # 标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "作品总结"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(1), Inches(8), Inches(0.4))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "PROJECT SUMMARY"
    p.font.size = Pt(14)
    p.font.color.rgb = C_CYAN
    
    page = s.shapes.add_textbox(Inches(12), Inches(0.35), Inches(1), Inches(0.5))
    tf = page.text_frame
    p = tf.paragraphs[0]
    p.text = "04"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 三个总结卡片
    summaries = [
        ("🎯", "技术选型", ["Next.js App Router", "React 19 + TypeScript", "Tailwind CSS", "Supabase 云数据库", "shadcn/ui 组件"], C_BLUE),
        ("✨", "创新特色", ["HR→Mi→Tcr递推算法", "实时数据可视化", "Web Bluetooth对接", "PMV-PPD热舒适度", "响应式设计"], C_PURPLE),
        ("🏆", "参赛信息", ["计算机设计大赛", "Web应用与开发赛道", "第十五届大赛", "现代Web技术栈", "团队合作开发"], C_ORANGE),
    ]
    
    x = 0.5
    for icon, title_text, points, color in summaries:
        # 卡片
        card = add_card(s, Inches(x), Inches(1.6), Inches(4), Inches(4.5), RGBColor(20, 40, 80), color)
        
        # 图标
        icon_box = s.shapes.add_textbox(Inches(x), Inches(1.8), Inches(4), Inches(1))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(3), Inches(3.6), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 要点
        y = 3.7
        for point in points:
            p_box = s.shapes.add_textbox(Inches(x + 0.4), Inches(y), Inches(3.2), Inches(0.45))
            tf = p_box.text_frame
            p = tf.paragraphs[0]
            p.text = "• " + point
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(180, 200, 220)
            y += 0.5
        
        x += 4.3
    
    # 底部感谢
    bottom = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.5), prs.slide_width, Inches(1))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = C_ORANGE
    bottom.line.fill.background()
    
    thanks = s.shapes.add_textbox(0, Inches(6.7), prs.slide_width, Inches(0.6))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听！欢迎交流！"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# ========== 生成 ==========
print("🎨 正在生成Web开发大赛PPT（多图版）...")

slide1_cover()
slide2_architecture()
slide3_highlights()
slide4_features()
slide5_summary()

output_path = "/workspace/projects/CDC健康检测_Web大赛PPT_多图版.pptx"
prs.save(output_path)
print(f"✅ PPT生成成功！")
print(f"📁 文件：{output_path}")
print(f"📊 共5页 - 包含4张技术图片")
