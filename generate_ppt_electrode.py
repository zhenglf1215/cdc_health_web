#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
有机表皮干电极产品介绍PPT（2页）
新加坡国立大学 - 人体心率和皮肤表面温度测量
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色 - 科技感
C_DARK = RGBColor(15, 25, 50)
C_BLUE = RGBColor(0, 120, 255)
C_CYAN = RGBColor(0, 220, 255)
C_ORANGE = RGBColor(255, 140, 0)
C_RED = RGBColor(255, 60, 60)
C_GREEN = RGBColor(0, 200, 120)
C_PURPLE = RGBColor(180, 100, 200)
C_WHITE = RGBColor(255, 255, 255)
C_LIGHT = RGBColor(248, 250, 255)
C_GRAY = RGBColor(100, 110, 120)

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_card(slide, left, top, width, height, fill=None, line_color=None):
    if fill is None:
        fill = C_WHITE
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if line_color:
        card.line.color.rgb = line_color
    else:
        card.line.fill.background()
    return card

# ========== 第1页：产品介绍 ==========
def slide1_product():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_DARK)
    
    # 装饰圆
    for x, y, size in [(Inches(-2), Inches(-2), Inches(6)), (Inches(10), Inches(4), Inches(6))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0, 80, 150)
        c.line.fill.background()
    
    # 机构标签
    tag = add_card(s, Inches(3), Inches(0.8), Inches(7.5), Inches(0.7), C_ORANGE)
    t = s.shapes.add_textbox(Inches(3), Inches(0.92), Inches(7.5), Inches(0.5))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "🏛️ 新加坡国立大学研究院"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 主标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.5), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "保形有机表皮干电极"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub = s.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(12.5), Inches(0.6))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "测量人体心率和皮肤表面温度"
    p.font.size = Pt(28)
    p.font.color.rgb = C_CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # 核心材料说明
    mat_title = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.5), Inches(0.5))
    tf = mat_title.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 核心材料：本征导电有机聚合物（PEDOT:PSS）+ 柔性弹性基底"
    p.font.size = Pt(16)
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # 四大特性
    features = [
        ("🤸", "机械柔性+自粘附", "完美贴合皮肤\n重物负载不脱落", C_BLUE),
        ("⚡", "高导电率", "稳定传输微弱信号\n可直接点亮LED", C_GREEN),
        ("💧", "抗汗可拉伸", "运动出汗稳定接触\n无凝胶失效问题", C_CYAN),
        ("🔗", "保形接触", "紧密贴合曲面\n降低接触阻抗", C_PURPLE),
    ]
    
    positions = [(0.5, 4.5), (3.5, 4.5), (6.5, 4.5), (9.5, 4.5)]
    
    for i, (icon, title_text, desc, color) in enumerate(features):
        x, y = positions[i]
        
        card = add_card(s, Inches(x), Inches(y), Inches(2.9), Inches(2.2), RGBColor(20, 40, 80), color)
        
        icon_box = s.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(2.9), Inches(0.8))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER
        
        t_box = s.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1), Inches(2.7), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        d_box = s.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.45), Inches(2.7), Inches(0.7))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(180, 200, 220)
        p.alignment = PP_ALIGN.CENTER

# ========== 第2页：测量原理 ==========
def slide2_principle():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_LIGHT)
    
    # 标题栏
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = C_DARK
    header.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "测量原理"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    # 左侧：心率测量
    hr_title = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.5))
    tf = hr_title.text_frame
    p = tf.paragraphs[0]
    p.text = "❤️ 心率测量原理（ECG心电图法）"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_RED
    
    hr_steps = [
        ("1", "电极贴附", "有机干电极粘贴胸部/手臂\n无需导电凝胶"),
        ("2", "信号采集", "捕捉皮肤表面微弱心电信号\n传输到采集模块"),
        ("3", "信号处理", "放大、滤波去除噪声\n还原心电图波形"),
        ("4", "心率计算", "识别QRS波群间隔\n计算每分钟心跳次数"),
    ]
    
    y = 1.8
    for num, step_title, desc in hr_steps:
        card = add_card(s, Inches(0.5), Inches(y), Inches(5.8), Inches(1.1), C_WHITE, C_RED)
        
        # 编号
        num_box = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y + 0.3), Inches(0.5), Inches(0.5))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = C_RED
        num_box.line.fill.background()
        n = s.shapes.add_textbox(Inches(0.7), Inches(y + 0.38), Inches(0.5), Inches(0.4))
        tf = n.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(1.4), Inches(y + 0.15), Inches(2), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        
        # 描述
        d_box = s.shapes.add_textbox(Inches(1.4), Inches(y + 0.55), Inches(4.7), Inches(0.5))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = C_GRAY
        
        y += 1.2
    
    # 右侧：温度测量
    temp_title = s.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(5), Inches(0.5))
    tf = temp_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🌡️ 皮肤温度测量原理（NTC热敏）"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_BLUE
    
    # 材料构成
    mat_card = add_card(s, Inches(6.8), Inches(1.8), Inches(6), Inches(1.5), RGBColor(240, 248, 255), C_BLUE)
    
    mat_title = s.shapes.add_textbox(Inches(7), Inches(1.95), Inches(5.6), Inches(0.4))
    tf = mat_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 热敏复合传感层材料构成"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_BLUE
    
    mats = [
        "导电相：PEDOT:PSS / 碳纳米管 / 石墨烯",
        "柔性基底：水性聚氨酯(WPU) / Ecoflex / PDMS",
        "增塑剂：山梨醇 / 离子液体",
    ]
    
    y = 2.4
    for mat in mats:
        m_box = s.shapes.add_textbox(Inches(7), Inches(y), Inches(5.6), Inches(0.3))
        tf = m_box.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + mat
        p.font.size = Pt(11)
        p.font.color.rgb = C_GRAY
        y += 0.32
    
    # 物理机制
    phy_card = add_card(s, Inches(6.8), Inches(3.5), Inches(6), Inches(2.2), RGBColor(240, 248, 255), C_BLUE)
    
    phy_title = s.shapes.add_textbox(Inches(7), Inches(3.65), Inches(5.6), Inches(0.4))
    tf = phy_title.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ 温度传感物理机制（NTC效应）"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_BLUE
    
    mechanisms = [
        ("🌡️", "温度升高 → 聚合物热膨胀\n→ 导电颗粒间距增大\n→ 导电通路概率下降\n→ 电阻上升", C_ORANGE),
        ("⚡", "温度升高 → 载流子迁移率提升\n→ 聚合物链段运动增强\n→ 电导率上升（主导NTC）\n→ 电阻下降", C_GREEN),
    ]
    
    y = 4.1
    for icon, desc, color in mechanisms:
        # 图标
        i_box = s.shapes.add_textbox(Inches(7), Inches(y), Inches(0.5), Inches(0.5))
        tf = i_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(18)
        
        # 描述
        d_box = s.shapes.add_textbox(Inches(7.5), Inches(y), Inches(5), Inches(0.9))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(9)
        p.font.color.rgb = C_GRAY
        p.line_spacing = 1.2
        
        y += 0.95
    
    # 底部特点
    bottom_card = add_card(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1), RGBColor(255, 245, 235))
    
    bottom_title = s.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.4))
    tf = bottom_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✨ 技术亮点：同基底共集成设计 + 功能解耦 + 力学匹配 + 工艺兼容"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    bottom_desc = s.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(11.5), Inches(0.4))
    tf = bottom_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "一片超薄可拉伸聚合物基底 = ECG干电极阵列 + 热敏复合传感层"
    p.font.size = Pt(12)
    p.font.color.rgb = C_GRAY
    p.alignment = PP_ALIGN.CENTER

# ========== 生成 ==========
print("🎨 正在生成有机表皮干电极产品介绍PPT...")

slide1_product()
slide2_principle()

output_path = "/workspace/projects/有机表皮干电极_产品介绍PPT.pptx"
prs.save(output_path)
print(f"✅ PPT生成成功！")
print(f"📁 文件：{output_path}")
print(f"📊 共2页 - 产品介绍 + 测量原理")
