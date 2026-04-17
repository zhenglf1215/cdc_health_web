#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CDC健康检测系统 - Web开发大赛PPT
聚焦：技术架构 + 开发创新点 + 前端亮点
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 16:9比例
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案 - 科技感
C_DARK = RGBColor(15, 25, 50)      # 深邃蓝黑
C_BLUE = RGBColor(0, 120, 255)     # 主蓝
C_CYAN = RGBColor(0, 220, 255)     # 亮青
C_ORANGE = RGBColor(255, 140, 0)   # 橙色
C_PURPLE = RGBColor(150, 80, 255)  # 紫色
C_GREEN = RGBColor(0, 200, 120)    # 绿色
C_WHITE = RGBColor(255, 255, 255)
C_LIGHT = RGBColor(248, 250, 255)
C_GRAY = RGBColor(100, 110, 120)

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_card(slide, left, top, width, height, fill=C_WHITE, line_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if line_color:
        card.line.color.rgb = line_color
    else:
        card.line.fill.background()
    return card

def add_tech_icon(slide, left, top, size, color, icon_text, name):
    """技术图标"""
    # 圆形背景
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    
    # 图标
    icon = slide.shapes.add_textbox(left, top + size/6, size, size/2)
    tf = icon.text_frame
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(int(size.pt * 0.4))
    p.alignment = PP_ALIGN.CENTER
    
    # 名称
    name_box = slide.shapes.add_textbox(left - Inches(0.3), top + size + Inches(0.1), size + Inches(0.6), Inches(0.4))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)
    p.font.color.rgb = C_GRAY
    p.alignment = PP_ALIGN.CENTER

# ========== 第1页：封面 ==========
def slide1_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_DARK)
    
    # 装饰圆
    for x, y, size, alpha in [(Inches(-1), Inches(-2), Inches(5), 30), 
                               (Inches(9), Inches(4), Inches(6), 25),
                               (Inches(11), Inches(-1), Inches(3), 15)]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BLUE
        c.line.fill.background()
    
    # 比赛标签
    tag = add_card(s, Inches(4.5), Inches(1.5), Inches(4.5), Inches(0.6), C_ORANGE)
    tag_text = s.shapes.add_textbox(Inches(4.5), Inches(1.6), Inches(4.5), Inches(0.5))
    tf = tag_text.text_frame
    p = tf.paragraphs[0]
    p.text = "🏆 计算机设计大赛"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 主标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.5), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "CDC健康检测系统"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.5), Inches(0.6))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Web前端开发作品"
    p.font.size = Pt(28)
    p.font.color.rgb = C_CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # 技术栈标签
    techs = ["Next.js 16", "React 19", "Tailwind CSS", "Supabase", "TypeScript"]
    x = 2.5
    for tech in techs:
        tag = add_card(s, Inches(x), Inches(4.8), Inches(1.7), Inches(0.45), RGBColor(0, 100, 180))
        t = s.shapes.add_textbox(Inches(x), Inches(4.88), Inches(1.7), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(12)
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        x += 1.9
    
    # 底部信息
    info = s.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.5), Inches(0.8))
    tf = info.text_frame
    p = tf.paragraphs[0]
    p.text = "第十五届大学生计算机设计大赛 | Web应用与开发赛道"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 170, 200)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "基于现代Web技术的健康监测平台"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(120, 140, 170)
    p.alignment = PP_ALIGN.CENTER

# ========== 第2页：技术架构 ==========
def slide2_architecture():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_LIGHT)
    
    # 标题栏
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = C_DARK
    header.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "技术架构"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    page = s.shapes.add_textbox(Inches(12), Inches(0.25), Inches(1), Inches(0.5))
    tf = page.text_frame
    p = tf.paragraphs[0]
    p.text = "01"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 架构层级 - 从上到下
    layers = [
        ("🌐", "用户端", ["Web浏览器", "响应式设计", "PWA支持"], C_BLUE),
        ("⚙️", "前端框架", ["Next.js 16", "React 19", "TypeScript"], C_PURPLE),
        ("🔌", "API层", ["Next.js API Routes", "RESTful设计", "数据校验"], C_CYAN),
        ("☁️", "云服务", ["Supabase", "PostgreSQL", "实时订阅"], C_GREEN),
        ("📡", "设备层", ["Web Bluetooth", "心率带对接", "实时采集"], C_ORANGE),
    ]
    
    y = 1.3
    for icon, name, desc_list, color in layers:
        # 层级卡片
        layer = add_card(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.0), C_WHITE, color)
        
        # 左侧色条
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.1), Inches(1.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # 图标
        icon_box = s.shapes.add_textbox(Inches(0.7), Inches(y + 0.2), Inches(0.8), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        
        # 名称
        name_box = s.shapes.add_textbox(Inches(1.5), Inches(y + 0.25), Inches(2), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        
        # 技术点
        x = 4
        for desc in desc_list:
            tag = add_card(s, Inches(x), Inches(y + 0.25), Inches(2.2), Inches(0.5), RGBColor(240, 245, 255))
            t = s.shapes.add_textbox(Inches(x), Inches(y + 0.35), Inches(2.2), Inches(0.4))
            tf = t.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(12)
            p.font.color.rgb = C_BLUE
            p.alignment = PP_ALIGN.CENTER
            x += 2.5
        
        y += 1.12

# ========== 第3页：开发亮点 ==========
def slide3_highlights():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_DARK)
    
    # 装饰
    for x, y, size in [(Inches(-2), Inches(-1), Inches(4)), (Inches(11), Inches(5), Inches(4))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0, 80, 160)
        c.line.fill.background()
    
    # 标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "前端开发亮点"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(1), Inches(8), Inches(0.4))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "FRONTEND HIGHLIGHTS"
    p.font.size = Pt(14)
    p.font.color.rgb = C_CYAN
    
    page = s.shapes.add_textbox(Inches(12), Inches(0.35), Inches(1), Inches(0.5))
    tf = page.text_frame
    p = tf.paragraphs[0]
    p.text = "02"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 六大亮点
    highlights = [
        ("⚡", "Next.js 16", "App Router\n服务端渲染\nSEO友好", C_BLUE),
        ("🎨", "Tailwind CSS", "原子化CSS\n响应式设计\n主题定制", C_PURPLE),
        ("🧩", "shadcn/ui", "Radix UI组件\nTypeScript\n无障碍支持", C_CYAN),
        ("📊", "数据可视化", "实时图表\nECharts集成\n交互体验", C_GREEN),
        ("🔄", "状态管理", "React Context\n数据持久化\n本地存储", C_ORANGE),
        ("🔒", "安全设计", "RLS权限控制\n数据校验\nXSS防护", RGBColor(255, 80, 120)),
    ]
    
    positions = [(0.4, 1.6), (4.5, 1.6), (8.6, 1.6),
                (0.4, 4.5), (4.5, 4.5), (8.6, 4.5)]
    
    for i, (icon, title_text, desc, color) in enumerate(highlights):
        x, y = positions[i]
        
        # 卡片
        card = add_card(s, Inches(x), Inches(y), Inches(4), Inches(2.7), RGBColor(20, 40, 80), color)
        
        # 图标
        icon_box = s.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(4), Inches(1))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.3), Inches(3.6), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        d_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.9), Inches(3.6), Inches(1))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(180, 200, 220)
        p.alignment = PP_ALIGN.CENTER

# ========== 第4页：核心功能展示 ==========
def slide4_features():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_LIGHT)
    
    # 标题栏
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = C_DARK
    header.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "亮点功能展示"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    page = s.shapes.add_textbox(Inches(12), Inches(0.25), Inches(1), Inches(0.5))
    tf = page.text_frame
    p = tf.paragraphs[0]
    p.text = "03"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 功能模块
    features = [
        ("👤", "用户管理", "登录注册\n角色权限\n头像系统", C_BLUE),
        ("📈", "数据看板", "实时心率\n趋势图表\n数据筛选", C_GREEN),
        ("🌡️", "体征监测", "多参数显示\n实时更新\n预警提醒", C_ORANGE),
        ("📱", "响应式设计", "PC/平板/手机\n自适应布局\n统一体验", C_PURPLE),
    ]
    
    positions = [(0.5, 1.2), (3.5, 1.2), (6.5, 1.2), (9.5, 1.2)]
    
    for i, (icon, title_text, desc, color) in enumerate(features):
        x, y = positions[i]
        
        # 卡片
        card = add_card(s, Inches(x), Inches(y), Inches(2.9), Inches(4.5), C_WHITE, color)
        
        # 顶部色条
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.9), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # 图标
        icon_box = s.shapes.add_textbox(Inches(x), Inches(y + 0.4), Inches(2.9), Inches(1.2))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(56)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.8), Inches(2.5), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        d_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 2.5), Inches(2.5), Inches(2))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = C_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.5
    
    # 底部技术标签
    tags = ["React Hooks", "Context API", "CSS Grid", "Flexbox", "Dark Mode", "Animation"]
    x = 0.8
    for tag in tags:
        t = add_card(s, Inches(x), Inches(6.0), Inches(2), Inches(0.45), RGBColor(240, 245, 255))
        text = s.shapes.add_textbox(Inches(x), Inches(6.1), Inches(2), Inches(0.35))
        tf = text.text_frame
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(11)
        p.font.color.rgb = C_BLUE
        p.alignment = PP_ALIGN.CENTER
        x += 2.2

# ========== 第5页：总结 ==========
def slide5_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C_DARK)
    
    # 装饰
    for x, y, size in [(Inches(-1), Inches(-1), Inches(4)), (Inches(10), Inches(5), Inches(5))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BLUE
        c.line.fill.background()
    
    # 标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "作品总结"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(1), Inches(8), Inches(0.4))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "PROJECT SUMMARY"
    p.font.size = Pt(14)
    p.font.color.rgb = C_CYAN
    
    page = s.shapes.add_textbox(Inches(12), Inches(0.35), Inches(1), Inches(0.5))
    tf = page.text_frame
    p = tf.paragraphs[0]
    p.text = "04"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 三个总结卡片
    summaries = [
        ("🎯", "技术选型", ["Next.js App Router", "React 19 + TypeScript", "Tailwind CSS + shadcn/ui", "Supabase 云数据库"], C_BLUE),
        ("✨", "创新特色", ["HR→Mi→Tcr递推算法", "实时数据可视化", "Web Bluetooth对接", "PMV-PPD热舒适度"], C_PURPLE),
        ("🏆", "参赛信息", ["计算机设计大赛", "Web应用与开发赛道", "第十五届大赛", "团队合作开发"], C_ORANGE),
    ]
    
    x = 0.5
    for icon, title_text, points, color in summaries:
        # 卡片
        card = add_card(s, Inches(x), Inches(1.6), Inches(4), Inches(4.5), RGBColor(20, 40, 80), color)
        
        # 图标
        icon_box = s.shapes.add_textbox(Inches(x), Inches(1.8), Inches(4), Inches(1))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(3), Inches(3.6), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 要点
        y = 3.7
        for point in points:
            p_box = s.shapes.add_textbox(Inches(x + 0.4), Inches(y), Inches(3.2), Inches(0.5))
            tf = p_box.text_frame
            p = tf.paragraphs[0]
            p.text = "• " + point
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(180, 200, 220)
            y += 0.5
        
        x += 4.3
    
    # 底部感谢
    bottom = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.5), prs.slide_width, Inches(1))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = C_ORANGE
    bottom.line.fill.background()
    
    thanks = s.shapes.add_textbox(0, Inches(6.7), prs.slide_width, Inches(0.6))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听！欢迎交流！"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# ========== 生成 ==========
print("🎨 正在生成Web开发大赛PPT...")

slide1_cover()
slide2_architecture()
slide3_highlights()
slide4_features()
slide5_summary()

output_path = "/workspace/projects/CDC健康检测_Web开发大赛PPT.pptx"
prs.save(output_path)
print(f"✅ PPT生成成功！")
print(f"📁 文件：{output_path}")
print(f"📊 共5页 - 聚焦Web开发技术与创新点")
