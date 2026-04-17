#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CDC健康检测系统 - 挑战杯PPT（高颜值多图版）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import math

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色方案 - 科技蓝+橙色点缀
C_DARK = RGBColor(15, 40, 80)      # 深蓝
C_BLUE = RGBColor(0, 120, 220)      # 主蓝
C_CYAN = RGBColor(0, 200, 255)      # 亮青
C_ORANGE = RGBColor(255, 140, 0)    # 橙色
C_WHITE = RGBColor(255, 255, 255)
C_LIGHT = RGBColor(245, 248, 255)   # 浅蓝白
C_GRAY = RGBColor(100, 110, 120)    # 灰色
C_GREEN = RGBColor(0, 200, 100)     # 绿色
C_RED = RGBColor(255, 60, 60)       # 红色

def add_gradient_shape(slide, left, top, width, height, color1, color2):
    """添加渐变效果的形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()
    return shape

def add_icon_circle(slide, left, top, size, color, icon_text):
    """添加图标圆形"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    
    # 添加图标文字
    icon_box = slide.shapes.add_textbox(left, top + size/4, size, size/2)
    tf = icon_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(int(size.pt * 0.4))
    p.alignment = PP_ALIGN.CENTER
    return circle

def add_arrow(slide, left, top, width, color):
    """添加箭头"""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.15))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow

# ========== 第1页：封面 ==========
def add_cover_slide(prs):
    slide = prs.slide_layouts[6]
    s = prs.slides.add_slide(slide)
    
    # 深色背景
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK
    bg.line.fill.background()
    
    # 装饰圆形
    for i, (x, y, size, alpha) in enumerate([
        (Inches(-1), Inches(-1), Inches(4), 40),
        (Inches(10), Inches(5), Inches(5), 30),
        (Inches(8), Inches(-2), Inches(3), 20),
    ]):
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = C_BLUE
        circle.line.fill.background()
    
    # 左侧装饰线
    line1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2), Inches(0.08), Inches(1.5))
    line1.fill.solid()
    line1.fill.fore_color.rgb = C_ORANGE
    line1.line.fill.background()
    
    # 主标题
    title = s.shapes.add_textbox(Inches(1.1), Inches(2), Inches(10), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "CDC健康检测系统"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    # 副标题
    sub = s.shapes.add_textbox(Inches(1.1), Inches(3.1), Inches(10), Inches(0.6))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "基于智能穿戴设备的热舒适度实时监测平台"
    p.font.size = Pt(24)
    p.font.color.rgb = C_CYAN
    
    # 底部信息卡片
    card_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.5), Inches(6), Inches(1.2))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RGBColor(0, 100, 180)
    card_bg.line.fill.background()
    
    info = s.shapes.add_textbox(Inches(1.3), Inches(5.7), Inches(5.5), Inches(1))
    tf = info.text_frame
    p = tf.paragraphs[0]
    p.text = "第十五届挑战杯"
    p.font.size = Pt(20)
    p.font.color.rgb = C_WHITE
    p = tf.add_paragraph()
    p.text = "大学生创业计划竞赛"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(180, 220, 255)
    
    # 右侧图标装饰
    icons = [("❤️", Inches(9), Inches(2.5)), ("🔥", Inches(10.5), Inches(3)), ("🌡️", Inches(9.5), Inches(4.5))]
    for icon, x, y in icons:
        icon_box = s.shapes.add_textbox(x, y, Inches(1.5), Inches(1))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
    
    # 底部条
    bottom = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), prs.slide_width, Inches(0.4))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = C_ORANGE
    bottom.line.fill.background()
    
    return s

# ========== 第2页：背景问题 ==========
def add_problem_slide(prs):
    slide = prs.slide_layouts[6]
    s = prs.slides.add_slide(slide)
    
    # 浅色背景
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_LIGHT
    bg.line.fill.background()
    
    # 顶部色条
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    top.fill.solid()
    top.fill.fore_color.rgb = C_DARK
    top.line.fill.background()
    
    # 标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(5), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "项目背景"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    # 三大问题 - 大图标卡片
    problems = [
        ("😷", "热射病频发", "高温作业死亡事故频发\n工人健康无法保障"),
        ("📊", "检测落后", "传统方式依赖人工\n效率低、误差大"),
        ("🔗", "数据孤岛", "心率体温代谢分散\n缺乏综合分析"),
    ]
    
    x = 0.8
    for icon, title_text, desc in problems:
        # 卡片
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2), Inches(3.8), Inches(4))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.fill.background()
        
        # 红色警告条
        warning = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.2), Inches(3.8), Inches(0.15))
        warning.fill.solid()
        warning.fill.fore_color.rgb = C_RED
        warning.line.fill.background()
        
        # 图标
        icon_box = s.shapes.add_textbox(Inches(x), Inches(1.6), Inches(3.8), Inches(1.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(72)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(3.3), Inches(3.4), Inches(0.6))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        d_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(4), Inches(3.4), Inches(1))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = C_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        x += 4.2
    
    # 箭头指向解决方案
    arrow_box = s.shapes.add_textbox(Inches(5.5), Inches(5.4), Inches(2), Inches(0.5))
    tf = arrow_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⬇"
    p.font.size = Pt(32)
    p.font.color.rgb = C_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # 底部解决方案预览
    solution_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.8), Inches(11), Inches(1.4))
    solution_bg.fill.solid()
    solution_bg.fill.fore_color.rgb = C_BLUE
    solution_bg.line.fill.background()
    
    sol_text = s.shapes.add_textbox(Inches(1.2), Inches(6.1), Inches(10.6), Inches(1))
    tf = sol_text.text_frame
    p = tf.paragraphs[0]
    p.text = "🚀 我们的方案：一站式智能监测平台"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return s

# ========== 第3页：核心功能 ==========
def add_function_slide(prs):
    slide = prs.slide_layouts[6]
    s = prs.slides.add_slide(slide)
    
    # 背景
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_LIGHT
    bg.line.fill.background()
    
    # 顶部
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    top.fill.solid()
    top.fill.fore_color.rgb = C_DARK
    top.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(5), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心功能"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    # 六个功能 - 2行3列大图标
    functions = [
        ("❤️", "心率监测", "实时采集", C_RED),
        ("🔥", "代谢率", "Mi计算", C_ORANGE),
        ("🌡️", "核心体温", "Tcr预测", RGBColor(200, 80, 80)),
        ("📊", "皮肤温度", "多点监测", C_GREEN),
        ("🏠", "CDC评估", "PMV-PPD", RGBColor(100, 100, 220)),
        ("👔", "服装建议", "智能推荐", RGBColor(180, 100, 200)),
    ]
    
    positions = [(0.5, 1.2), (4.6, 1.2), (8.7, 1.2),
                 (0.5, 4.2), (4.6, 4.2), (8.7, 4.2)]
    
    for i, (icon, title_text, sub, color) in enumerate(functions):
        x, y = positions[i]
        
        # 卡片
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.fill.background()
        
        # 顶部色条
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # 大图标
        icon_box = s.shapes.add_textbox(Inches(x), Inches(y + 0.3), Inches(3.8), Inches(1.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(64)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.9), Inches(3.4), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        s_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 2.4), Inches(3.4), Inches(0.4))
        tf = s_box.text_frame
        p = tf.paragraphs[0]
        p.text = sub
        p.font.size = Pt(14)
        p.font.color.rgb = C_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return s

# ========== 第4页：创新点（重点页）==========
def add_innovation_slide(prs):
    slide = prs.slide_layouts[6]
    s = prs.slides.add_slide(slide)
    
    # 深色背景
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK
    bg.line.fill.background()
    
    # 装饰圆
    for x, y, size in [(Inches(-2), Inches(-1), Inches(4)), (Inches(11), Inches(5), Inches(4))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BLUE
        c.line.fill.background()
    
    # 标题
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心创新点"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    # 副标题
    sub = s.shapes.add_textbox(Inches(0.5), Inches(1), Inches(6), Inches(0.4))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "CORE INNOVATION"
    p.font.size = Pt(14)
    p.font.color.rgb = C_CYAN
    
    # 三大创新点 - 左右布局
    innovations = [
        ("💡", "01", "HR→Mi→Tcr\n递推算法", "通过心率变化\n预测核心体温\n提前15-30分钟预警", C_ORANGE),
        ("📈", "02", "多参数融合\nPMV-PPD", "综合环境+代谢率\n+体温数据\n精准舒适度评估", C_BLUE),
        ("⚠️", "03", "三级预警\n响应机制", "🔵正常 🔶注意 🔴危险\n智能分级告警\n保障作业安全", C_RED),
    ]
    
    x = 0.5
    for icon, num, title_text, desc, color in innovations:
        # 卡片
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(4), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(20, 50, 100)
        card.line.color.rgb = color
        
        # 编号
        num_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(1.8), Inches(1), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        
        # 图标
        icon_box = s.shapes.add_textbox(Inches(x + 0.3), Inches(2.5), Inches(3.4), Inches(1.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(56)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        t_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(4.2), Inches(3.6), Inches(1))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        d_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(5.3), Inches(3.6), Inches(1.4))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(180, 200, 220)
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.3
        
        x += 4.3
    
    # 底部公式条
    formula_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7), prs.slide_width, Inches(0.5))
    formula_bg.fill.solid()
    formula_bg.fill.fore_color.rgb = RGBColor(0, 80, 160)
    formula_bg.line.fill.background()
    
    return s

# ========== 第5页：总结展望 ==========
def add_summary_slide(prs):
    slide = prs.slide_layouts[6]
    s = prs.slides.add_slide(slide)
    
    # 背景
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_LIGHT
    bg.line.fill.background()
    
    # 顶部
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    top.fill.solid()
    top.fill.fore_color.rgb = C_DARK
    top.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(5), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "总结与展望"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    # 左侧成果
    card1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(4), Inches(3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = C_WHITE
    card1.line.fill.background()
    
    # 成果图标
    icons = ["✅", "✅", "✅", "✅", "✅"]
    results = ["完整Web监测系统", "HR→Mi→Tcr算法", "PMV-PPD评估", "蓝牙设备对接", "双端界面"]
    
    for i, (icon, text) in enumerate(zip(icons, results)):
        icon_box = s.shapes.add_textbox(Inches(0.7), Inches(1.3 + i*0.5), Inches(0.5), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(16)
        
        t_box = s.shapes.add_textbox(Inches(1.2), Inches(1.35 + i*0.5), Inches(3), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(15)
        p.font.color.rgb = C_DARK
    
    # 中间应用场景
    card2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.1), Inches(4), Inches(3))
    card2.fill.solid()
    card2.fill.fore_color.rgb = C_WHITE
    card2.line.fill.background()
    
    scenes = [
        ("🏗️", "高温作业", "建筑/钢铁/煤矿"),
        ("🏃", "运动员", "马拉松/足球"),
        ("🏥", "医疗", "术后监测"),
        ("👴", "养老", "异常预警"),
    ]
    
    for i, (icon, title, sub) in enumerate(scenes):
        icon_box = s.shapes.add_textbox(Inches(5), Inches(1.3 + i*0.7), Inches(0.8), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(24)
        
        t_box = s.shapes.add_textbox(Inches(5.8), Inches(1.35 + i*0.7), Inches(2), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        
        s_box = s.shapes.add_textbox(Inches(7.5), Inches(1.4 + i*0.7), Inches(1.2), Inches(0.4))
        tf = s_box.text_frame
        p = tf.paragraphs[0]
        p.text = sub
        p.font.size = Pt(11)
        p.font.color.rgb = C_GRAY
    
    # 右侧未来规划
    card3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(1.1), Inches(3.8), Inches(3))
    card3.fill.solid()
    card3.fill.fore_color.rgb = RGBColor(255, 245, 235)
    card3.line.color.rgb = C_ORANGE
    
    future_title = s.shapes.add_textbox(Inches(9.3), Inches(1.2), Inches(3.4), Inches(0.4))
    tf = future_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🚀 未来规划"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    
    futures = ["小程序端开发", "多设备支持", "AI预警模型", "工业级部署"]
    for i, f in enumerate(futures):
        arrow = s.shapes.add_textbox(Inches(9.3), Inches(1.7 + i*0.55), Inches(0.4), Inches(0.4))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(14)
        p.font.color.rgb = C_ORANGE
        
        t_box = s.shapes.add_textbox(Inches(9.7), Inches(1.75 + i*0.55), Inches(3), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = f
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
    
    # 底部技术指标
    metrics_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.4), Inches(12.4), Inches(1.8))
    metrics_bg.fill.solid()
    metrics_bg.fill.fore_color.rgb = C_WHITE
    metrics_bg.line.fill.background()
    
    metrics = [
        ("1次/秒", "心率采样"),
        ("<500ms", "数据延迟"),
        (">95%", "Tcr准确率"),
        ("<200ms", "系统响应"),
    ]
    
    x = 1.5
    for value, label in metrics:
        # 数值
        v_box = s.shapes.add_textbox(Inches(x), Inches(4.6), Inches(2.5), Inches(0.8))
        tf = v_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = C_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # 标签
        l_box = s.shapes.add_textbox(Inches(x), Inches(5.5), Inches(2.5), Inches(0.4))
        tf = l_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = C_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        x += 3
    
    # 底部感谢
    thanks_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.5), prs.slide_width, Inches(1))
    thanks_bg.fill.solid()
    thanks_bg.fill.fore_color.rgb = C_DARK
    thanks_bg.line.fill.background()
    
    thanks = s.shapes.add_textbox(0, Inches(6.7), prs.slide_width, Inches(0.8))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听  欢迎交流"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return s

# ========== 生成PPT ==========
print("🎨 正在生成高颜值PPT...")

add_cover_slide(prs)
add_problem_slide(prs)
add_function_slide(prs)
add_innovation_slide(prs)
add_summary_slide(prs)

output_path = "/workspace/projects/CDC健康检测_挑战杯PPT_多图版.pptx"
prs.save(output_path)
print(f"✅ PPT生成成功！")
print(f"📁 文件：{output_path}")
print(f"📊 共5页 - 图片多、文字少、高颜值")
