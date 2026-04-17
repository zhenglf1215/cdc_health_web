#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CDC健康检测系统 - 挑战杯PPT生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
PRIMARY_BLUE = RgbColor(0, 123, 255)
DARK_BLUE = RgbColor(0, 45, 114)
ORANGE = RgbColor(255, 140, 0)
RED = RgbColor(220, 53, 69)
GREEN = RgbColor(40, 167, 69)
WHITE = RgbColor(255, 255, 255)
DARK_GRAY = RgbColor(52, 58, 64)
LIGHT_GRAY = RgbColor(233, 236, 239)

def add_title_slide(prs):
    """第1页：封面"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色块 - 渐变效果（用多个矩形模拟）
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(240, 248, 255)
    shape.line.fill.background()
    
    # 顶部蓝色条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.fill.background()
    
    # 底部蓝色条
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(6.5), prs.slide_width, Inches(1)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = DARK_BLUE
    bottom_bar.line.fill.background()
    
    # 左侧装饰条
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(1.5), Inches(0.3), Inches(5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = ORANGE
    left_bar.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CDC 健康检测系统"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "基于智能穿戴设备的热舒适度实时监测平台"
    p.font.size = Pt(28)
    p.font.color.rgb = PRIMARY_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.3), Inches(5), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()
    
    # 比赛信息
    info_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11), Inches(1.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "第十五届挑战杯大学生创业计划竞赛"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "XX大学 XX团队 | 2025年4月"
    p.font.size = Pt(20)
    p.font.color.rgb = RgbColor(108, 117, 125)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_background_slide(prs):
    """第2页：项目背景与痛点"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(250, 250, 252)
    bg.line.fill.background()
    
    # 顶部标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "项目背景与痛点"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 页码
    page_box = slide.shapes.add_textbox(Inches(12), Inches(0.35), Inches(1), Inches(0.5))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = "01"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 左侧 - 行业痛点
    left_label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.5), Inches(2.5), Inches(0.5)
    )
    left_label.fill.solid()
    left_label.fill.fore_color.rgb = RED
    left_label.line.fill.background()
    tf = left_label.text_frame
    tf.paragraphs[0].text = "行业痛点"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False
    
    # 痛点卡片
    pain_points = [
        ("热射病频发", "我国每年因高温作业导致热射病死亡数十人，建筑工人、钢铁工人、煤矿工人是高危群体"),
        ("检测手段落后", "传统CDC检测依赖人工问询、查表计算，效率低、误差大、无法实时监测"),
        ("数据孤岛严重", "心率、体温、代谢数据分散在不同设备，缺乏综合分析和预警能力")
    ]
    
    y_pos = 2.2
    for title, desc in pain_points:
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y_pos), Inches(5.8), Inches(1.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RgbColor(220, 220, 220)
        
        # 红色标记
        mark = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.6), Inches(y_pos + 0.15), Inches(0.25), Inches(0.25)
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = RED
        mark.line.fill.background()
        
        # 标题
        t_box = slide.shapes.add_textbox(Inches(1.0), Inches(y_pos + 0.1), Inches(5), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RED
        
        # 描述
        d_box = slide.shapes.add_textbox(Inches(1.0), Inches(y_pos + 0.55), Inches(5), Inches(0.8))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        
        y_pos += 1.55
    
    # 右侧 - 我们的方案
    right_label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(2.5), Inches(0.5)
    )
    right_label.fill.solid()
    right_label.fill.fore_color.rgb = GREEN
    right_label.line.fill.background()
    tf = right_label.text_frame
    tf.paragraphs[0].text = "我们的方案"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 方案卡片
    solutions = [
        ("一体化监测平台", "Web + 蓝牙设备 + 云端数据库，心率、体温、代谢、舒适度一站式监测"),
        ("智能化算法支撑", "基于生理参数的递推预测模型，提前预警热应激风险"),
        ("实时数据可视化", "图表直观展示历史趋势，支持多人多设备管理")
    ]
    
    y_pos = 2.2
    for title, desc in solutions:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(y_pos), Inches(5.8), Inches(1.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RgbColor(220, 220, 220)
        
        mark = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(7.0), Inches(y_pos + 0.15), Inches(0.25), Inches(0.25)
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = GREEN
        mark.line.fill.background()
        
        t_box = slide.shapes.add_textbox(Inches(7.4), Inches(y_pos + 0.1), Inches(5), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GREEN
        
        d_box = slide.shapes.add_textbox(Inches(7.4), Inches(y_pos + 0.55), Inches(5), Inches(0.8))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        
        y_pos += 1.55
    
    return slide

def add_functions_slide(prs):
    """第3页：核心功能"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(250, 250, 252)
    bg.line.fill.background()
    
    # 顶部标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心功能模块"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    page_box = slide.shapes.add_textbox(Inches(12), Inches(0.35), Inches(1), Inches(0.5))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = "02"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 功能模块
    functions = [
        ("❤️", "心率监测", "实时采集心率数据\n蓝牙心率带秒级更新\n历史趋势分析", PRIMARY_BLUE),
        ("🔥", "劳动代谢率", "Mi能量消耗计算\n基于心率推算\n运动强度评估", ORANGE),
        ("🌡️", "核心体温", "Tcr递推算法预测\n热风险预警\n37.5℃告警", RED),
        ("📊", "皮肤温度", "多点温度监测\n热图可视化展示\n历史曲线分析", RgbColor(40, 167, 69)),
        ("🏠", "CDC评估", "PMV-PPD热舒适度\n国际标准模型\n作业建议", RgbColor(111, 66, 193)),
        ("👔", "服装建议", "个性化着装推荐\n适应环境温度\n降低热负荷", RgbColor(253, 126, 20))
    ]
    
    positions = [
        (0.4, 1.5), (4.5, 1.5), (8.6, 1.5),
        (0.4, 4.2), (4.5, 4.2), (8.6, 4.2)
    ]
    
    for i, (icon, title, desc, color) in enumerate(functions):
        x, y = positions[i]
        
        # 卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RgbColor(230, 230, 230)
        
        # 顶部色条
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = color
        top_bar.line.fill.background()
        
        # 图标和标题
        icon_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.25), Inches(0.8), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        
        title_box = slide.shapes.add_textbox(Inches(x + 1.0), Inches(y + 0.3), Inches(2.5), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.9), Inches(3.4), Inches(1.4))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.3
    
    return slide

def add_innovation_slide(prs):
    """第4页：技术架构与创新点（重点页）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(250, 250, 252)
    bg.line.fill.background()
    
    # 顶部标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "技术架构与核心创新点 ⭐"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    page_box = slide.shapes.add_textbox(Inches(12), Inches(0.35), Inches(1), Inches(0.5))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = "03"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 左侧 - 技术架构
    arch_label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.4), Inches(2), Inches(0.45)
    )
    arch_label.fill.solid()
    arch_label.fill.fore_color.rgb = PRIMARY_BLUE
    arch_label.line.fill.background()
    tf = arch_label.text_frame
    tf.paragraphs[0].text = "技术架构"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 架构层
    layers = [
        ("Web 网页端", "React + Tailwind CSS\n响应式设计", RgbColor(52, 152, 219)),
        ("Next.js 16", "SSR + API Routes\nTypeScript类型安全", RgbColor(155, 89, 182)),
        ("Supabase 云数据库", "PostgreSQL\n实时同步 + RLS权限", RgbColor(46, 204, 113)),
        ("蓝牙心率带", "Web Bluetooth API\n跨平台兼容", RgbColor(231, 76, 60))
    ]
    
    y = 2.0
    for name, desc, color in layers:
        # 层背景
        layer = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(5.8), Inches(1.0)
        )
        layer.fill.solid()
        layer.fill.fore_color.rgb = WHITE
        layer.line.color.rgb = color
        layer.line.width = Pt(2)
        
        # 左侧色块
        left = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(y), Inches(0.2), Inches(1.0)
        )
        left.fill.solid()
        left.fill.fore_color.rgb = color
        left.line.fill.background()
        
        # 名称
        n_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(3), Inches(0.4))
        tf = n_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        # 描述
        d_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.5), Inches(5), Inches(0.5))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        
        y += 1.15
    
    # 右侧 - 核心创新点
    inno_label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(1.4), Inches(2.5), Inches(0.45)
    )
    inno_label.fill.solid()
    inno_label.fill.fore_color.rgb = ORANGE
    inno_label.line.fill.background()
    tf = inno_label.text_frame
    tf.paragraphs[0].text = "核心创新点 ⭐"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 创新点1
    inno1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(2.0), Inches(6.3), Inches(1.6)
    )
    inno1.fill.solid()
    inno1.fill.fore_color.rgb = RgbColor(255, 248, 240)
    inno1.line.color.rgb = ORANGE
    
    # 创新点1标题
    t1 = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(6), Inches(0.4))
    tf = t1.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 创新点1：HR → Mi → Tcr 递推算法"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    # 创新点1内容
    c1 = slide.shapes.add_textbox(Inches(6.8), Inches(2.5), Inches(6), Inches(1.1))
    tf = c1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• 传统：仅测量皮肤表面温度\n• 创新：通过心率变化预测核心体温变化\n• 效果：提前15-30分钟预警热射病风险"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.line_spacing = 1.3
    
    # 创新点2
    inno2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(3.75), Inches(6.3), Inches(1.6)
    )
    inno2.fill.solid()
    inno2.fill.fore_color.rgb = RgbColor(240, 248, 255)
    inno2.line.color.rgb = PRIMARY_BLUE
    
    t2 = slide.shapes.add_textbox(Inches(6.8), Inches(3.85), Inches(6), Inches(0.4))
    tf = t2.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 创新点2：多参数融合的PMV-PPD模型"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    
    c2 = slide.shapes.add_textbox(Inches(6.8), Inches(4.25), Inches(6), Inches(1.1))
    tf = c2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• 综合：环境温度+湿度+风速+辐射+服装热阻\n• 融合：代谢率(Mi)+核心体温(Tcr)\n• 优势：比单一指标更准确的舒适度评估"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.line_spacing = 1.3
    
    # 创新点3
    inno3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(5.5), Inches(6.3), Inches(1.6)
    )
    inno3.fill.solid()
    inno3.fill.fore_color.rgb = RgbColor(255, 240, 240)
    inno3.line.color.rgb = RED
    
    t3 = slide.shapes.add_textbox(Inches(6.8), Inches(5.6), Inches(6), Inches(0.4))
    tf = t3.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 创新点3：三级预警响应机制"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RED
    
    c3 = slide.shapes.add_textbox(Inches(6.8), Inches(6.0), Inches(6), Inches(1.0))
    tf = c3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• 🔵 蓝色 = 正常    • 🔶 橙色 = 注意    • 🔴 红色 = 危险\n• Tsk>36.5℃ → 橙色  |  Tcr>37.5℃ → 橙色  |  Tcr>38.0℃ → 红色"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.line_spacing = 1.4
    
    # 公式展示
    formula_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.7)
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = RgbColor(243, 244, 246)
    formula_box.line.fill.background()
    
    f_text = slide.shapes.add_textbox(Inches(0.6), Inches(6.75), Inches(12), Inches(0.5))
    tf = f_text.text_frame
    p = tf.paragraphs[0]
    p.text = "核心公式：Mi = 65 + (HR-HRrest)/(180-0.65×Age-HRrest) × [(41.7-0.22×Age)×W^(2/3)-65]    Tcr(t+1) = Tcr(t) + 0.0036×(Mi-55)×0.0952"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_summary_slide(prs):
    """第5页：总结与展望"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(250, 250, 252)
    bg.line.fill.background()
    
    # 顶部标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "总结与展望"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    page_box = slide.shapes.add_textbox(Inches(12), Inches(0.35), Inches(1), Inches(0.5))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = "04"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.RIGHT
    
    # 左侧 - 项目成果
    result_label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.4), Inches(2), Inches(0.45)
    )
    result_label.fill.solid()
    result_label.fill.fore_color.rgb = GREEN
    result_label.line.fill.background()
    tf = result_label.text_frame
    tf.paragraphs[0].text = "项目成果 ✅"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    results = [
        "完整的CDC健康监测Web系统",
        "HR→Mi→Tcr核心算法实现",
        "PMV-PPD热舒适度评估",
        "Web Bluetooth设备对接",
        "管理端+用户端双界面"
    ]
    
    y = 2.0
    for result in results:
        check = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(0.4), Inches(0.4))
        tf = check.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(16)
        p.font.color.rgb = GREEN
        p.font.bold = True
        
        t_box = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(5), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = result
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        y += 0.5
    
    # 技术指标
    metric_label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.6), Inches(2), Inches(0.4)
    )
    metric_label.fill.solid()
    metric_label.fill.fore_color.rgb = PRIMARY_BLUE
    metric_label.line.fill.background()
    tf = metric_label.text_frame
    tf.paragraphs[0].text = "技术指标"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    metrics = [
        ("心率采样", "1次/秒"),
        ("数据延迟", "<500ms"),
        ("Tcr准确率", ">95%"),
        ("系统响应", "<200ms")
    ]
    
    y = 5.15
    for name, value in metrics:
        n_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2.5), Inches(0.35))
        tf = n_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        
        v_box = slide.shapes.add_textbox(Inches(2.5), Inches(y), Inches(1.5), Inches(0.35))
        tf = v_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE
        y += 0.38
    
    # 中间 - 应用场景
    scene_label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(1.4), Inches(2), Inches(0.45)
    )
    scene_label.fill.solid()
    scene_label.fill.fore_color.rgb = RgbColor(111, 66, 193)
    scene_label.line.fill.background()
    tf = scene_label.text_frame
    tf.paragraphs[0].text = "应用场景"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    scenes = [
        ("🏗️", "高温作业", "建筑、钢铁、煤矿"),
        ("🏃", "运动员训练", "马拉松、足球"),
        ("🏥", "医疗康复", "术后体温监测"),
        ("👴", "养老院", "老人异常预警")
    ]
    
    y = 2.0
    for icon, title, desc in scenes:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(y), Inches(4.2), Inches(0.75)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RgbColor(220, 220, 220)
        
        icon_box = slide.shapes.add_textbox(Inches(4.4), Inches(y + 0.15), Inches(0.6), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(22)
        
        t_box = slide.shapes.add_textbox(Inches(5.1), Inches(y + 0.1), Inches(3), Inches(0.35))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        d_box = slide.shapes.add_textbox(Inches(5.1), Inches(y + 0.4), Inches(3), Inches(0.35))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        
        y += 0.85
    
    # 右侧 - 未来规划
    future_label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.4), Inches(2), Inches(0.45)
    )
    future_label.fill.solid()
    future_label.fill.fore_color.rgb = ORANGE
    future_label.line.fill.background()
    tf = future_label.text_frame
    tf.paragraphs[0].text = "未来规划 🚀"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    futures = [
        ("小程序端开发", "便捷移动端访问"),
        ("多设备支持", "血氧、血压、睡眠监测"),
        ("AI预警模型", "机器学习风险预测"),
        ("工业级部署", "企业级解决方案")
    ]
    
    y = 2.0
    for title, desc in futures:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(y), Inches(4.1), Inches(0.75)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RgbColor(220, 220, 220)
        
        # 箭头
        arrow = slide.shapes.add_textbox(Inches(9.0), Inches(y + 0.2), Inches(0.4), Inches(0.4))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        
        t_box = slide.shapes.add_textbox(Inches(9.5), Inches(y + 0.1), Inches(3), Inches(0.35))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        d_box = slide.shapes.add_textbox(Inches(9.5), Inches(y + 0.4), Inches(3), Inches(0.35))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        
        y += 0.85
    
    # 底部结语
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(6.8), prs.slide_width, Inches(0.7)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = DARK_BLUE
    bottom_bar.line.fill.background()
    
    end_box = slide.shapes.add_textbox(0, Inches(6.95), prs.slide_width, Inches(0.5))
    tf = end_box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢聆听！欢迎提问与交流！"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# 生成PPT
print("正在生成PPT...")
add_title_slide(prs)
add_background_slide(prs)
add_functions_slide(prs)
add_innovation_slide(prs)
add_summary_slide(prs)

# 保存
output_path = "/workspace/projects/CDC健康检测系统_挑战杯PPT.pptx"
prs.save(output_path)
print(f"✅ PPT生成成功！")
print(f"📁 文件位置：{output_path}")
print(f"📊 共 5 页幻灯片")
